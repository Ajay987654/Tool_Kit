import os
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
from reportlab.pdfgen import canvas

def convert_file(filepath, target):
    filename, ext = os.path.splitext(filepath)
    output_file = ""

    if ext.lower() == ".pdf" and target == "word":
        doc = Document()
        reader = PdfReader(filepath)
        for page in reader.pages:
            text = page.extract_text()
            if text:
                doc.add_paragraph(text)
        output_file = filename + ".docx"
        doc.save(output_file)


    elif ext.lower() == ".pdf" and target == "ppt":
        prs = Presentation()
        reader = PdfReader(filepath)
        for page in reader.pages:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            text = page.extract_text() or " "
            slide.shapes.title.text = text[:50]
            if len(text) > 50:
                slide.placeholders[1].text = text[50:]
        output_file = filename + ".pptx"
        prs.save(output_file)

    elif ext.lower() == ".docx" and target == "pdf":
        doc = Document(filepath)
        output_file = filename + ".pdf"
        c = canvas.Canvas(output_file)
        y = 800
        for para in doc.paragraphs:
            c.drawString(50, y, para.text)
            y -= 20
        c.save()

    elif ext.lower() == ".docx" and target == "ppt":
        prs = Presentation()
        doc = Document(filepath)
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = text[:50]
            if len(text) > 50:
                slide.placeholders[1].text = text[50:]
        output_file = filename + ".pptx"
        prs.save(output_file)

    else:
        raise Exception("Unsupported conversion")

    return output_file
